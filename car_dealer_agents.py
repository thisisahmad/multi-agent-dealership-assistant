# car_dealer_agents.py
# Streamlined 6-agent car dealership system using OpenAI Agents SDK
# 6 Core Agents: Orchestrator, Temporal, Knowledge, Booking, Handoff, Lead Capture

import os
import glob
import json
import re
import argparse
import logging
import asyncio
from datetime import datetime, timedelta
from typing import Optional, Dict, Any

import numpy as np
import faiss
import pdfplumber
from pptx import Presentation
from dotenv import load_dotenv
from openai import OpenAI

try:
    from agents.agent import Agent
    from agents.run import Runner
    from agents.tool import function_tool
    from agents.model_settings import ModelSettings
    from agents.memory import SQLiteSession
except ImportError:
    # Fallback for older SDK versions
    from agents import (
        Agent,
        Runner,
        function_tool,
        ModelSettings,
        SQLiteSession,
    )

# =================== SETUP ===================

load_dotenv()
logging.getLogger("pdfminer").setLevel(logging.ERROR)
logging.getLogger("pdfplumber").setLevel(logging.ERROR)
logging.getLogger("reportlab").setLevel(logging.ERROR)

api_key = os.getenv("OPENAI_API_KEY")
if not api_key:
    raise ValueError("Please set OPENAI_API_KEY in .env as OPENAI_API_KEY")

client = OpenAI(api_key=api_key)

EMBED_MODEL = "text-embedding-3-small"
FAST_MODEL = os.getenv("OPENAI_FAST_MODEL", "gpt-4o-mini")
GEN_MODEL = os.getenv("OPENAI_MODEL", FAST_MODEL)
TEMPORAL_MODEL = os.getenv("OPENAI_TEMPORAL_MODEL", FAST_MODEL)
CHUNK_SIZE = 2600
OVERLAP = 300
TOP_K = 3  # Reduced for faster search

AGENT_LOG: list[dict] = []  # collected per-run

# Performance optimization: In-memory cache for FAISS indices and chunks
_FAISS_CACHE: dict[str, tuple] = {}  # {dealer_name: (index, chunks)}
_EMBEDDING_CACHE: dict[str, np.ndarray] = {}  # {text: embedding}
_CONFIG_CACHE: dict[str, dict] = {}
_FORM_CACHE: dict[str, dict] = {}

# Dealer context for RAG
_CURRENT_DEALER_NAME: str = None

# Form state management
_FORM_STATES: dict[str, dict] = {}  # {session_id: {form_id, fields, current_field}}
_HANDOFF_STATE: dict[str, dict] = {}  # {session_id: {awaiting_details: bool}}


def load_config(config_name: str) -> dict:
    """Load JSON config file with in-memory caching."""
    if config_name in _CONFIG_CACHE:
        return _CONFIG_CACHE[config_name]

    config_path = os.path.join("configs", config_name)
    if os.path.exists(config_path):
        with open(config_path, "r", encoding="utf-8") as f:
            _CONFIG_CACHE[config_name] = json.load(f)
            return _CONFIG_CACHE[config_name]

    _CONFIG_CACHE[config_name] = {}
    return _CONFIG_CACHE[config_name]


def log_event(agent: str, message: str):
    rec = {"time": datetime.utcnow().isoformat(), "agent": agent, "message": message}
    AGENT_LOG.append(rec)
    print(f"[{agent}] {message}")


# =================== PROMPTS ===================

languages = {
    "ar": "Arabic", "bg": "Bulgarian", "hr": "Croatian", "cs": "Czech",
    "dk": "Danish", "nl": "Dutch", "en": "English", "fi": "Finnish",
    "fr": "French", "it": "Italian", "gr": "Greek", "de": "German",
    "hu": "Hungarian", "jp": "Japanese", "no": "Norwegian", "pl": "Polish",
    "pt": "Portuguese", "ro": "Romanian", "sk": "Slovak", "sl": "Slovenian",
    "es": "Spanish", "se": "Swedish",
}

HANDOVER = """
#### Handover to Human Agent Guidelines

Avoid suggesting customers contact the dealership or customer support; _you_ are the primary contact
for the dealership, but can transfer the conversation to a human agent.

Recognize the importance of a human touch and readily recommend a human representative when necessary.
The option to speak to a human is always available to the customer.

Initiate a handover to a human representative when:
1. The customer requests to switch to a human agent.
2. The bot faces a sophisticated customer question beyond its capability.
3. The topic is related to a critical situation, sensitive information, or security.
4. The customer expresses dissatisfaction or uses emotion-laden keywords.
5. The user explicitly asks to speak with a human.

When in doubt, ask the user if they would like to be connected to a human agent.
"""


# =================== VECTOR KB (RAG) ===================

def get_embedding(text: str) -> np.ndarray:
    """Get embedding with caching for performance."""
    if text in _EMBEDDING_CACHE:
        return _EMBEDDING_CACHE[text]
    
    r = client.embeddings.create(model=EMBED_MODEL, input=text)
    embedding = np.array(r.data[0].embedding, dtype="float32")
    
    if len(_EMBEDDING_CACHE) < 1000:
        _EMBEDDING_CACHE[text] = embedding
    
    return embedding


def get_embeddings_batch(texts: list[str]) -> list[list[float]]:
    r = client.embeddings.create(model=EMBED_MODEL, input=texts)
    return [d.embedding for d in r.data]


def extract_text_from_pdf(path: str) -> str:
    txt = ""
    try:
        import fitz
        with fitz.open(path) as doc:
            for page in doc:
                t = page.get_text("text")
                if t:
                    txt += t + "\n"
        return txt
    except Exception:
        pass

    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            t = page.extract_text()
            if t:
                txt += t + "\n"
    return txt


def extract_text_from_ppt(path: str) -> str:
    prs = Presentation(path)
    txt = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                txt += shape.text + "\n"
    return txt


def chunk_text(s: str, size=CHUNK_SIZE, overlap=OVERLAP) -> list[str]:
    out = []
    i, n = 0, len(s)
    while i < n:
        j = i + size
        chunk = s[i:j].strip()
        if chunk:
            out.append(chunk)
        i = j - overlap
    return out


def detect_dealer_from_filename(filename: str) -> Optional[str]:
    """
    Detect dealer/company name from filename. Generalized - works with any company name.
    
    Examples:
    - "Sytner_20_QA.pdf" → "Sytner"
    - "Lookers_info.pdf" → "Lookers"
    - "Charles_Hurst_20_QA.pdf" → "Charles_Hurst"
    - "CompanyName_document.pdf" → "CompanyName"
    """
    import re
    name = os.path.splitext(filename)[0]
    
    # Pattern 1: CompanyName_* (e.g., "Sytner_20_QA", "Charles_Hurst_20_QA")
    # Extract up to 2 words separated by underscore (for multi-word company names)
    match = re.match(r'^([A-Za-z]+(?:_[A-Za-z]+)?)', name)
    if match:
        dealer = match.group(1)
        # Filter out common non-company prefixes (generalized list)
        excluded = ['data', 'kb', 'updated', 'car', 'file', 'document', 'test', 'temp']
        if dealer.lower() not in excluded:
            return dealer
    
    # Pattern 2: CompanyName-* (e.g., "Lookers-Service")
    match = re.match(r'^([A-Za-z]+)', name)
    if match:
        dealer = match.group(1)
        excluded = ['data', 'kb', 'updated', 'car', 'file', 'document', 'test', 'temp']
        if dealer.lower() not in excluded:
            return dealer
    
    return None


def get_available_dealers(data_dir="Data"):
    """Get list of available dealers from filenames in Data/."""
    dealers = set()
    if not os.path.exists(data_dir):
        return []
    
    for item in os.listdir(data_dir):
        item_path = os.path.join(data_dir, item)
        if os.path.isfile(item_path):
            if item.endswith(('.index', '_chunks.json')):
                continue
            dealer = detect_dealer_from_filename(item)
            if dealer:
                dealers.add(dealer)
    
    return sorted(list(dealers))


def build_kb_from_data(dealer_name: str = None, data_dir="Data", index_path=None, docs_path=None, auto_detect_dealers: bool = True):
    """
    Build KB from documents. Files stay in Data/, KB files saved in Data/{dealer}/.
    
    If dealer_name is None and auto_detect_dealers=True:
    - Auto-detects all dealers from filenames
    - Builds separate KB for each dealer
    - Generalizes to any dealer name found in filenames
    
    If dealer_name is provided:
    - Builds KB only for that specific dealer
    """
    if dealer_name:
        # Build KB for specific dealer
        if index_path is None:
            kb_dir = os.path.join(data_dir, dealer_name)
            os.makedirs(kb_dir, exist_ok=True)
            index_path = os.path.join(kb_dir, "kb.index")
        if docs_path is None:
            kb_dir = os.path.join(data_dir, dealer_name)
            os.makedirs(kb_dir, exist_ok=True)
            docs_path = os.path.join(kb_dir, "kb_chunks.json")
        
        files = []
        if os.path.exists(data_dir):
            for item in os.listdir(data_dir):
                item_path = os.path.join(data_dir, item)
                if os.path.isfile(item_path):
                    if item.endswith(('.index', '_chunks.json')):
                        continue
                    detected_dealer = detect_dealer_from_filename(item)
                    if detected_dealer and detected_dealer.lower() == dealer_name.lower():
                        files.append(item_path)
        
        print(f"🏢 Building KB for dealer: {dealer_name}")
        print(f"📄 Found {len(files)} files matching {dealer_name}")
        
        if not files:
            raise ValueError(f"No files found for dealer '{dealer_name}' in {data_dir}/")
        
        chunks: list[str] = []
        for p in files:
            print(f"🗂️ Processing: {os.path.basename(p)}")
            if p.lower().endswith(".pdf"):
                t = extract_text_from_pdf(p)
            elif p.lower().endswith((".pptx", ".ppt")):
                t = extract_text_from_ppt(p)
            else:
                with open(p, "r", encoding="utf-8") as f:
                    t = f.read()
            chunks.extend(chunk_text(t))
        
        print(f"🔧 Embedding {len(chunks)} chunks...")
        embs = get_embeddings_batch(chunks)
        dim = len(embs[0])
        idx = faiss.IndexFlatL2(dim)
        idx.add(np.array(embs, dtype="float32"))
        faiss.write_index(idx, index_path)
        with open(docs_path, "w", encoding="utf-8") as f:
            json.dump(chunks, f, ensure_ascii=False)
        
        print(f"✅ KB built for {dealer_name}: {len(chunks)} chunks → {index_path}")
        
        # Clear cache after rebuild
        if dealer_name in _FAISS_CACHE:
            del _FAISS_CACHE[dealer_name]
    
    else:
        # Auto-detect dealers and build separate KBs for each
        if auto_detect_dealers:
            dealers = get_available_dealers(data_dir)
            
            if dealers:
                # Build separate KB for each detected dealer
                print(f"🔍 Auto-detected {len(dealers)} dealers: {', '.join(dealers)}")
                print(f"🏗️ Building separate KB files for each dealer...\n")
                
                for dealer in dealers:
                    print(f"\n{'='*50}")
                    print(f"📦 Processing dealer: {dealer}")
                    print(f"{'='*50}")
                    try:
                        build_kb_from_data(dealer_name=dealer, data_dir=data_dir, auto_detect_dealers=False)
                    except Exception as e:
                        print(f"❌ Error building KB for {dealer}: {e}")
                        continue
                
                print(f"\n{'='*50}")
                print(f"✅ Completed building KBs for all dealers")
                print(f"{'='*50}")
                return
        
        # No dealers detected or auto_detect_dealers=False, build general KB
        dealers = get_available_dealers(data_dir) if auto_detect_dealers else []
        if not dealers:
            # No dealers detected, build general KB
            if index_path is None:
                index_path = "kb.index"
            if docs_path is None:
                docs_path = "kb_chunks.json"
            
            files = glob.glob(os.path.join(data_dir, "*"))
            files = [f for f in files if os.path.isfile(f) and not f.endswith(('.index', '_chunks.json'))]
            
            if not files:
                raise ValueError(f"No files found in {data_dir}/")
            
            print(f"🏢 Building general KB (no dealers detected)")
            chunks: list[str] = []
            for p in files:
                print(f"🗂️ Processing: {os.path.basename(p)}")
                if p.lower().endswith(".pdf"):
                    t = extract_text_from_pdf(p)
                elif p.lower().endswith((".pptx", ".ppt")):
                    t = extract_text_from_ppt(p)
                else:
                    with open(p, "r", encoding="utf-8") as f:
                        t = f.read()
                chunks.extend(chunk_text(t))
            
            print(f"🔧 Embedding {len(chunks)} chunks...")
            embs = get_embeddings_batch(chunks)
            dim = len(embs[0])
            idx = faiss.IndexFlatL2(dim)
            idx.add(np.array(embs, dtype="float32"))
            faiss.write_index(idx, index_path)
            with open(docs_path, "w", encoding="utf-8") as f:
                json.dump(chunks, f, ensure_ascii=False)
            
            print(f"✅ KB built: {len(chunks)} chunks → {index_path}")
            
            if "general" in _FAISS_CACHE:
                del _FAISS_CACHE["general"]
        else:
            # Build separate KB for each detected dealer
            print(f"🔍 Auto-detected {len(dealers)} dealers: {', '.join(dealers)}")
            print(f"🏗️ Building separate KB files for each dealer...\n")
            
            for dealer in dealers:
                print(f"\n{'='*50}")
                print(f"📦 Processing dealer: {dealer}")
                print(f"{'='*50}")
                try:
                    build_kb_from_data(dealer_name=dealer, data_dir=data_dir, auto_detect_dealers=False)
                except Exception as e:
                    print(f"❌ Error building KB for {dealer}: {e}")
                    continue
            
            print(f"\n{'='*50}")
            print(f"✅ Completed building KBs for all dealers")
            print(f"{'='*50}")


def ensure_kb(dealer_name: str = None, index_path=None, docs_path=None):
    """Ensure KB exists. KB files stored in Data/{dealer_name}/ folder."""
    if dealer_name:
        if index_path is None:
            kb_dir = os.path.join("Data", dealer_name)
            index_path = os.path.join(kb_dir, "kb.index")
        if docs_path is None:
            kb_dir = os.path.join("Data", dealer_name)
            docs_path = os.path.join(kb_dir, "kb_chunks.json")
    else:
        if index_path is None:
            index_path = "kb.index"
        if docs_path is None:
            docs_path = "kb_chunks.json"
    
    if not (os.path.exists(index_path) and os.path.exists(docs_path)):
        build_kb_from_data(dealer_name=dealer_name, index_path=index_path, docs_path=docs_path)


def _load_kb_to_cache(dealer_name: str = None):
    """Load KB into memory cache for faster access."""
    cache_key = dealer_name or "general"
    
    if cache_key in _FAISS_CACHE:
        return _FAISS_CACHE[cache_key]
    
    if dealer_name:
        kb_dir = os.path.join("Data", dealer_name)
        index_path = os.path.join(kb_dir, "kb.index")
        docs_path = os.path.join(kb_dir, "kb_chunks.json")
    else:
        index_path = "kb.index"
        docs_path = "kb_chunks.json"
    
    ensure_kb(dealer_name=dealer_name, index_path=index_path, docs_path=docs_path)
    
    idx = faiss.read_index(index_path)
    with open(docs_path, "r", encoding="utf-8") as f:
        docs = json.load(f)
    
    _FAISS_CACHE[cache_key] = (idx, docs)
    return idx, docs


def set_dealer_name(dealer_name: str):
    """Set the current dealer name for RAG operations."""
    global _CURRENT_DEALER_NAME
    _CURRENT_DEALER_NAME = dealer_name

def get_dealer_name() -> str:
    """Get the current dealer name."""
    return _CURRENT_DEALER_NAME


def detect_dealer_in_query(query: str) -> Optional[str]:
    """
    Detect if query mentions a dealer name. Dynamic - checks against available dealers.
    Returns dealer name if found, None otherwise.
    Handles variations: underscores, spaces, case-insensitive.
    """
    query_lower = query.lower().replace("_", " ").replace("-", " ")
    available_dealers = get_available_dealers("Data")
    
    for dealer in available_dealers:
        # Normalize dealer name (handle underscores, spaces)
        dealer_normalized = dealer.lower().replace("_", " ").replace("-", " ")
        dealer_words = dealer_normalized.split()
        
        # Check if all words of dealer name appear in query
        if all(word in query_lower for word in dealer_words):
            # Check if it's a word boundary match (not part of another word)
            # Try both original and normalized versions
            patterns = [
                r'\b' + re.escape(dealer.lower()) + r'\b',
                r'\b' + re.escape(dealer_normalized) + r'\b',
            ]
            for pattern in patterns:
                if re.search(pattern, query_lower):
                    return dealer
    
    return None


def check_dealer_mismatch_logic(query: str) -> Dict[str, Any]:
    """
    Check if query mentions a different dealer than currently selected.
    Returns dict with 'is_mismatch' (bool) and 'message' (str).
    Use this BEFORE calling semantic_search to prevent cross-dealer answers.
    """
    current_dealer = get_dealer_name()
    query_dealer = detect_dealer_in_query(query)
    
    if not current_dealer:
        return {
            "is_mismatch": False,
            "message": "No dealer selected. Using general KB.",
            "current_dealer": None,
            "query_dealer": query_dealer
        }
    
    if query_dealer and query_dealer.lower() != current_dealer.lower():
        return {
            "is_mismatch": True,
            "message": f"I don't have information about {query_dealer}. I can only provide information about {current_dealer}.",
            "current_dealer": current_dealer,
            "query_dealer": query_dealer
        }
    
    return {
        "is_mismatch": False,
        "message": f"Query is about current dealer '{current_dealer}' or no dealer mentioned.",
        "current_dealer": current_dealer,
        "query_dealer": query_dealer
    }


@function_tool
def check_dealer_mismatch(query: str) -> Dict[str, Any]:
    return check_dealer_mismatch_logic(query)


def semantic_search_logic(query: str, k: int = TOP_K) -> Dict[str, Any]:
    """
    RAG retrieval logic. Uses in-memory cache for faster performance.
    STRICT: Only searches the currently selected dealer's KB.
    """
    current_dealer = get_dealer_name()
    query_dealer = detect_dealer_in_query(query)

    if query_dealer and current_dealer and query_dealer.lower() != current_dealer.lower():
        log_event("semantic_search", f"BLOCKED: Query about {query_dealer} but dealer is {current_dealer}")
        return {
            "error": True,
            "blocked": True,
            "message": f"I don't have information about {query_dealer}. I can only provide information about {current_dealer}.",
            "current_dealer": current_dealer,
            "query_dealer": query_dealer,
            "results": [],
            "instruction": "You MUST respond with the exact message above. Do NOT generate any other answer.",
        }

    idx, docs = _load_kb_to_cache(current_dealer)
    q = get_embedding(query)
    D, I = idx.search(np.array([q]), k=k)
    results = []
    dealer_label = f" ({current_dealer})" if current_dealer else ""
    for rank, i in enumerate(I[0], 1):
        results.append(
            {"rank": rank, "text": docs[i], "source": f"Data{dealer_label}/* (chunk {i})"}
        )

    return {
        "error": False,
        "blocked": False,
        "message": f"Found {len(results)} results from {current_dealer if current_dealer else 'general'} KB",
        "results": results,
        "current_dealer": current_dealer,
    }


@function_tool
def semantic_search(query: str, k: int = TOP_K) -> Dict[str, Any]:
    """
    RAG retrieval tool. Uses in-memory cache for faster performance.
    STRICT: Only searches the currently selected dealer's KB.
    CRITICAL: You MUST call check_dealer_mismatch() FIRST before calling this tool.
    If query mentions a different dealer, returns error message.
    """
    return semantic_search_logic(query, k=k)


@function_tool
def retrieve_chunks(query: str, k: int = TOP_K) -> Dict[str, Any]:
    """
    Retrieve text chunks from KB using semantic search. 
    CRITICAL: You MUST call check_dealer_mismatch() FIRST before calling this tool.
    Returns dict with results or error message (same as semantic_search).
    """
    return semantic_search_logic(query, k)


def _not_in_kb_message() -> str:
    behaviour = load_config("behaviour_rules.json")
    return (
        behaviour.get("response_guidelines", {})
        .get("knowledge_agent", {})
        .get("if_not_in_kb", "This information is not available in the knowledge base.")
    )


def _normalize_for_match(text: str) -> list[str]:
    stop_words = {
        "a",
        "an",
        "and",
        "are",
        "about",
        "can",
        "do",
        "does",
        "for",
        "from",
        "how",
        "i",
        "is",
        "me",
        "of",
        "sell",
        "sells",
        "the",
        "their",
        "they",
        "to",
        "types",
        "what",
        "when",
        "where",
        "who",
        "with",
        "you",
        "your",
    }
    words = re.findall(r"[a-z0-9]+", text.lower())
    return [word for word in words if len(word) > 2 and word not in stop_words]


def _extract_direct_qa_answer(query: str, chunks: list[str]) -> Optional[str]:
    """
    The dealer PDFs are numbered Q&A documents. If a retrieved chunk contains
    the matching question, return its answer directly instead of asking the LLM
    to infer it from a large context block.
    """
    query_terms = set(_normalize_for_match(query))
    if not query_terms:
        return None

    best_score = 0
    best_answer = None
    qa_pattern = re.compile(
        r"(?:^|\n)\s*\d+\.\s*(?P<question>[^\n?]+\?)\s*\n"
        r"(?P<answer>.*?)(?=\n\s*\d+\.\s*[^\n?]+\?\s*\n|\Z)",
        re.DOTALL,
    )

    for chunk in chunks:
        for match in qa_pattern.finditer(chunk):
            question = match.group("question").strip()
            answer = re.sub(r"\s+", " ", match.group("answer")).strip()
            if not answer:
                continue

            question_terms = set(_normalize_for_match(question))
            overlap = query_terms.intersection(question_terms)
            score = len(overlap)

            if score > best_score:
                best_score = score
                best_answer = answer

    # Two meaningful shared terms is enough for Q&A titles like
    # "When was Charles Hurst founded?" or "What types of vehicles..."
    if best_score >= 2:
        return best_answer

    return None


def answer_knowledge_query_logic(query: str, dealer_name: str = None) -> str:
    """
    Deterministic RAG path used by the orchestrator for general enquiries.
    Searches the selected dealer KB and answers only from retrieved chunks.
    """
    not_in_kb = _not_in_kb_message()

    mismatch = check_dealer_mismatch_logic(query)
    if mismatch["is_mismatch"]:
        log_event("Knowledge", mismatch["message"])
        return mismatch["message"]

    search_result = semantic_search_logic(query)
    if search_result.get("blocked"):
        return search_result["message"]

    results = search_result.get("results", [])
    if not results:
        log_event("Knowledge", "No KB chunks retrieved")
        return not_in_kb

    retrieved_chunks = [result["text"] for result in results]
    direct_answer = _extract_direct_qa_answer(query, retrieved_chunks)
    if direct_answer:
        log_event("Knowledge", "Answered from direct Q&A match")
        return direct_answer

    context = "\n\n---\n\n".join(retrieved_chunks)
    dealer_label = dealer_name or search_result.get("current_dealer") or "the dealership"

    log_event("Knowledge", f"Retrieved {len(results)} chunks for query")

    response = client.chat.completions.create(
        model=FAST_MODEL,
        temperature=0.0,
        max_tokens=400,
        messages=[
            {
                "role": "system",
                "content": (
                    f"You answer questions about {dealer_label} using ONLY the provided context. "
                    f"If the answer is not clearly supported by the context, respond EXACTLY: "
                    f"\"{not_in_kb}\". Do not guess or use outside knowledge. Be concise."
                ),
            },
            {
                "role": "user",
                "content": f"Context:\n{context}\n\nQuestion: {query}",
            },
        ],
    )

    answer = response.choices[0].message.content.strip()
    log_event("Knowledge", f"Generated grounded answer ({len(answer)} chars)")
    return answer


def calculate_day_of_week_logic(date_str: str, tz: str = "Europe/London") -> Dict[str, Any]:
    """Calculate day of week for a date string. Handles relative dates like 'coming Thursday'."""
    try:
        import pytz
        now = datetime.now(pytz.timezone(tz))
    except Exception:
        now = datetime.utcnow()
    
    date_str_lower = date_str.lower().strip()
    today_day = now.strftime("%A")
    days_of_week = ["Monday", "Tuesday", "Wednesday", "Thursday", "Friday", "Saturday", "Sunday"]
    today_index = days_of_week.index(today_day)
    
    if date_str_lower == "today":
        target_date = now
    elif date_str_lower == "tomorrow":
        target_date = now + timedelta(days=1)
    elif date_str_lower.startswith(("coming ", "this ")):
        day_name = date_str_lower.replace("coming ", "").replace("this ", "").strip().capitalize()
        if day_name in days_of_week:
            target_index = days_of_week.index(day_name)
            days_ahead = (target_index - today_index) % 7
            if days_ahead == 0:
                days_ahead = 7
            target_date = now + timedelta(days=days_ahead)
        else:
            return {"valid": False, "error": f"Invalid day name: {day_name}"}
    elif date_str_lower.startswith("next "):
        day_name = date_str_lower.replace("next ", "").strip().capitalize()
        if day_name in days_of_week:
            target_index = days_of_week.index(day_name)
            days_ahead = (target_index - today_index) % 7
            if days_ahead == 0:
                days_ahead = 7
            target_date = now + timedelta(days=days_ahead + 7)
        else:
            return {"valid": False, "error": f"Invalid day name: {day_name}"}
    else:
        try:
            from dateutil import parser
            target_date = parser.parse(date_str, default=now)
        except Exception:
            return {"valid": False, "error": f"Could not parse date: {date_str}"}
    
    day_name = target_date.strftime("%A")
    full_date = target_date.strftime("%A, %d %B %Y")
    iso_date = target_date.strftime("%Y-%m-%d")
    is_weekend = day_name in ["Saturday", "Sunday"]
    days_from_today = (target_date.date() - now.date()).days
    
    return {
        "valid": True,
        "day_of_week": day_name,
        "full_date": full_date,
        "iso_date": iso_date,
        "is_weekend": is_weekend,
        "days_from_today": days_from_today,
    }


def _current_time_context(tz: str = "Europe/London") -> tuple[datetime, str]:
    try:
        import pytz

        now = datetime.now(pytz.timezone(tz))
    except Exception:
        now = datetime.utcnow()

    return now, now.strftime("%A, %d %B %Y %H:%M")


def _json_from_model(content: str) -> dict:
    cleaned = content.replace("```json", "").replace("```", "").strip()
    return json.loads(cleaned)


def normalize_human_date_logic(date_text: str, tz: str = "Europe/London") -> Dict[str, Any]:
    """Normalize natural-language dates like 'next week Sunday' or '20 June'."""
    now, now_label = _current_time_context(tz)

    def _normalize_date_local() -> Optional[Dict[str, Any]]:
        normalized = date_text.lower().strip()
        days_of_week = ["Monday", "Tuesday", "Wednesday", "Thursday", "Friday", "Saturday", "Sunday"]
        today_index = days_of_week.index(now.strftime("%A"))

        if normalized in {"today", "tomorrow"}:
            target = now if normalized == "today" else now + timedelta(days=1)
            return {
                "valid": True,
                "iso_date": target.strftime("%Y-%m-%d"),
                "full_date": target.strftime("%A, %d %B %Y"),
                "display_date": target.strftime("%A, %d %B %Y"),
            }

        for i, day in enumerate(days_of_week):
            if day.lower() in normalized:
                days_ahead = (i - today_index) % 7
                if "next week" in normalized:
                    days_ahead += 7 if days_ahead != 0 else 7
                elif normalized.startswith("next ") or days_ahead == 0:
                    days_ahead = days_ahead or 7

                target = now + timedelta(days=days_ahead)
                return {
                    "valid": True,
                    "iso_date": target.strftime("%Y-%m-%d"),
                    "full_date": target.strftime("%A, %d %B %Y"),
                    "display_date": target.strftime("%A, %d %B %Y"),
                }

        try:
            from dateutil import parser

            target = parser.parse(date_text, default=now)
            if target.date() < now.date():
                target = target.replace(year=target.year + 1)
            return {
                "valid": True,
                "iso_date": target.strftime("%Y-%m-%d"),
                "full_date": target.strftime("%A, %d %B %Y"),
                "display_date": target.strftime("%A, %d %B %Y"),
            }
        except Exception:
            return None

    local_result = _normalize_date_local()
    if local_result:
        return local_result

    try:
        response = client.chat.completions.create(
            model=TEMPORAL_MODEL,
            messages=[
                {
                    "role": "system",
                    "content": (
                        "You convert human date text into a concrete calendar date. "
                        "Use the provided current date/time and timezone. "
                        "Return JSON only with keys: valid, iso_date, display_date, reason. "
                        "If the text is ambiguous but still reasonably interpretable, choose the next future date."
                    ),
                },
                {
                    "role": "user",
                    "content": (
                        f"Current date/time: {now_label} ({tz})\n"
                        f"User date text: {date_text}"
                    ),
                },
            ],
        )
        data = _json_from_model(response.choices[0].message.content or "{}")
        if data.get("valid") and data.get("iso_date"):
            parsed = datetime.strptime(data["iso_date"], "%Y-%m-%d")
            return {
                "valid": True,
                "iso_date": data["iso_date"],
                "full_date": parsed.strftime("%A, %d %B %Y"),
                "display_date": data.get("display_date") or parsed.strftime("%A, %d %B %Y"),
            }
    except Exception as e:
        log_event("Temporal", f"GPT date normalization fallback: {e}")

    return {"valid": False, "error": f"I couldn't understand that date: {date_text}"}


def normalize_human_time_logic(time_text: str, tz: str = "Europe/London") -> Dict[str, Any]:
    """Normalize natural-language times like '8pm', 'half past two', or 'morning'."""
    now, now_label = _current_time_context(tz)

    def _normalize_time_local() -> Optional[Dict[str, Any]]:
        text = time_text.lower().strip().replace(".", "")
        direct = re.match(r"^(\d{1,2})(?::(\d{2}))?\s*(am|pm)?$", text)
        if direct:
            hour = int(direct.group(1))
            minute = int(direct.group(2) or 0)
            meridiem = direct.group(3)
            if meridiem == "pm" and hour < 12:
                hour += 12
            elif meridiem == "am" and hour == 12:
                hour = 0
            if 0 <= hour <= 23 and 0 <= minute <= 59:
                hour_12 = hour % 12 or 12
                suffix = "AM" if hour < 12 else "PM"
                return {
                    "valid": True,
                    "time_24h": f"{hour:02d}:{minute:02d}",
                    "display_time": f"{hour_12}:{minute:02d} {suffix}",
                }

        broad_times = {
            "morning": "10:00",
            "late morning": "11:00",
            "afternoon": "14:00",
            "evening": "17:00",
            "noon": "12:00",
            "midday": "12:00",
        }
        for phrase, normalized_time in broad_times.items():
            if phrase in text:
                return {
                    "valid": True,
                    "time_24h": normalized_time,
                    "display_time": normalized_time,
                }
        return None

    local_result = _normalize_time_local()
    if local_result:
        return local_result

    try:
        response = client.chat.completions.create(
            model=TEMPORAL_MODEL,
            messages=[
                {
                    "role": "system",
                    "content": (
                        "You convert human time text into 24-hour HH:MM time. "
                        "Return JSON only with keys: valid, time_24h, display_time, reason. "
                        "If the user gives a broad period like morning or afternoon, choose a sensible dealership appointment time."
                    ),
                },
                {
                    "role": "user",
                    "content": (
                        f"Current date/time: {now_label} ({tz})\n"
                        f"User time text: {time_text}"
                    ),
                },
            ],
        )
        data = _json_from_model(response.choices[0].message.content or "{}")
        if data.get("valid") and data.get("time_24h"):
            parsed = datetime.strptime(data["time_24h"], "%H:%M")
            return {
                "valid": True,
                "time_24h": parsed.strftime("%H:%M"),
                "display_time": data.get("display_time") or parsed.strftime("%H:%M"),
            }
    except Exception as e:
        log_event("Temporal", f"GPT time normalization fallback: {e}")

    return {"valid": False, "error": f"I couldn't understand that time: {time_text}"}


@function_tool
def calculate_day_of_week(date_str: str, tz: str = "Europe/London") -> Dict[str, Any]:
    return calculate_day_of_week_logic(date_str, tz)


@function_tool
def check_business_hours(date_str: str, time_str: Optional[str] = None, tz: str = "Europe/London") -> Dict[str, Any]:
    """Check if date/time is within business hours (Mon-Fri 9:00-17:00, Sat 9:00-13:00, Sun closed)."""
    result = calculate_day_of_week_logic(date_str, tz)
    if not result.get("valid"):
        return result
    
    day = result["day_of_week"]
    
    if day == "Sunday":
        return {
            "valid": False,
            "in_business_hours": False,
            "reason": "We are closed on Sundays",
            "day": day,
            "date": result["full_date"],
        }
    
    if time_str:
        try:
            hour, minute = map(int, time_str.split(":"))
            if day == "Saturday":
                in_hours = 9 <= hour < 13
                reason = "Saturday hours: 9:00-13:00" if not in_hours else "Within business hours"
            else:
                in_hours = 9 <= hour < 17
                reason = "Weekday hours: 9:00-17:00" if not in_hours else "Within business hours"
            
            return {
                "valid": True,
                "in_business_hours": in_hours,
                "reason": reason,
                "day": day,
                "date": result["full_date"],
                "time": time_str,
            }
        except Exception:
            return {"valid": False, "error": f"Invalid time format: {time_str}"}
    
    return {
        "valid": True,
        "in_business_hours": True,
        "day": day,
        "date": result["full_date"],
        "business_hours": "9:00-13:00" if day == "Saturday" else "9:00-17:00",
    }


@function_tool
def validate_appointment(date_str: str, time_slot: Optional[str] = None, appointment_type: str = "workshop") -> Dict[str, Any]:
    """Validate appointment date and time."""
    hours_check = check_business_hours(date_str, time_slot)
    
    if not hours_check.get("valid") or not hours_check.get("in_business_hours"):
        return {
            "available": False,
            "reason": hours_check.get("reason", "Outside business hours"),
            "suggested_alternatives": ["Try a weekday between 9:00-17:00"]
        }
    
    if time_slot:
        return {
            "available": True,
            "date": hours_check.get("date"),
            "time_slot": time_slot,
            "status": "available"
        }
    else:
        day = hours_check.get("day", "")
        if day == "Saturday":
            slots = ["9:00", "10:00", "11:00", "12:00"]
        else:
            slots = ["9:00", "10:00", "11:00", "14:00", "15:00", "16:00"]
        
        return {
            "available": True,
            "date": hours_check.get("date"),
            "suggested_slots": slots,
            "status": "available"
        }


@function_tool
def check_agent_availability() -> Dict[str, Any]:
    return check_agent_availability_logic()


def check_agent_availability_logic() -> Dict[str, Any]:
    """Check if human agents are currently available."""
    try:
        import pytz
        now = datetime.now(pytz.timezone("Europe/London"))
    except Exception:
        now = datetime.utcnow()
    
    hour = now.hour
    day = now.strftime("%A")
    
    if day == "Sunday":
        available = False
        reason = "Agents available Monday-Saturday"
    elif day == "Saturday":
        available = 9 <= hour < 13
        reason = "Limited hours on Saturday (9:00-13:00)" if not available else "Available"
    else:
        available = 9 <= hour < 17
        reason = "Outside business hours (9:00-17:00)" if not available else "Available"
    
    return {
        "available": available,
        "reason": reason,
        "current_time": now.strftime("%A %d %B %Y %H:%M"),
        "estimated_wait": "5-10 minutes" if available else "Next business day"
    }


@function_tool
def capture_lead(customer_name: str, email: Optional[str] = None, phone: Optional[str] = None,
                 interest_type: Optional[str] = None, vehicle_interest: Optional[str] = None) -> Dict[str, Any]:
    return capture_lead_logic(customer_name, email, phone, interest_type, vehicle_interest)


def capture_lead_logic(customer_name: str, email: Optional[str] = None, phone: Optional[str] = None,
                       interest_type: Optional[str] = None, vehicle_interest: Optional[str] = None) -> Dict[str, Any]:
    """Capture lead information for follow-up."""
    lead_data = {
        "customer_name": customer_name,
        "email": email,
        "phone": phone,
        "interest_type": interest_type or "general_enquiry",
        "vehicle_interest": vehicle_interest,
        "timestamp": datetime.utcnow().isoformat(),
        "status": "captured"
    }
    
    log_event("Lead Capture Agent", f"Lead captured: {customer_name} - {interest_type}")
    
    return {
        "success": True,
        "lead_id": f"LEAD-{datetime.utcnow().strftime('%Y%m%d%H%M%S')}",
        "message": "Lead information captured successfully. A team member will contact you shortly.",
        "data": lead_data
    }


@function_tool
def extract_entities(conversation_text: str, entity_types: Optional[list[str]] = None) -> Dict[str, Any]:
    return extract_entities_logic(conversation_text)


def extract_entities_logic(conversation_text: str) -> Dict[str, Any]:
    """Extract structured entities from conversation text."""
    entities = {}
    import re
    
    email_pattern = r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b'
    emails = re.findall(email_pattern, conversation_text)
    if emails:
        entities["email"] = emails[0]
    
    phone_pattern = r'(\+44|0)[0-9]{10,11}'
    phones = re.findall(phone_pattern, conversation_text)
    if phones:
        entities["phone"] = phones[0]
    
    return {
        "entities": entities,
        "count": len(entities),
        "extracted_at": datetime.utcnow().isoformat()
    }


@function_tool
def validate_contact_info(email: Optional[str] = None, phone: Optional[str] = None) -> Dict[str, Any]:
    """Validate contact information format."""
    import re
    errors = []
    validated = {}
    
    if email:
        email_pattern = r'^[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}$'
        if re.match(email_pattern, email):
            validated["email"] = email
        else:
            errors.append("Invalid email format")
    
    if phone:
        phone_pattern = r'^(\+44|0)[0-9]{10,11}$'
        if re.match(phone_pattern, phone.replace(" ", "").replace("-", "")):
            validated["phone"] = phone
        else:
            errors.append("Invalid phone format (expected UK format)")
    
    return {
        "valid": len(errors) == 0,
        "validated": validated,
        "errors": errors
    }


@function_tool
def classify_intent(user_message: str, conversation_context: str = "") -> Dict[str, Any]:
    return classify_intent_logic(user_message)


def classify_intent_logic(user_message: str) -> Dict[str, Any]:
    """Classify user intent from message. Returns intent_id and description."""
    intents_config = load_config("intents.json")
    intents = intents_config.get("intents", [])
    
    message_lower = user_message.lower()
    
    # Score-based keyword matching (returns intent with highest score)
    best_intent = None
    best_score = 0
    
    for intent in intents:
        keywords = intent.get("keywords", [])
        score = 0
        matched_keywords = []
        
        for keyword in keywords:
            if keyword.lower() in message_lower:
                # Longer keywords get higher scores (more specific)
                keyword_score = len(keyword.split())
                score += keyword_score
                matched_keywords.append(keyword)
        
        if score > best_score:
            best_score = score
            best_intent = intent
    
    if best_intent:
        return {
            "intent_id": best_intent["intent_id"],
            "description": best_intent["description"],
            "form": best_intent.get("form"),
            "confidence": min(0.9, 0.5 + (best_score * 0.1))
        }
    
    # Default to general enquiry
    return {
        "intent_id": "general_enquiry",
        "description": "General information request",
        "form": None,
        "confidence": 0.5
    }


def get_form_id_from_intent_logic(intent_id: str) -> Dict[str, Any]:
    """Get form_id from intent_id. Returns form_id to use for booking."""
    intents_config = load_config("intents.json")
    intents = intents_config.get("intents", [])
    
    for intent in intents:
        if intent.get("intent_id") == intent_id:
            form_id = intent.get("form")
            if form_id:
                return {"form_id": form_id, "intent_id": intent_id}
    
    # Default to service_booking for booking-related intents
    if intent_id in ["service_booking", "test_drive"]:
        return {"form_id": intent_id, "intent_id": intent_id}
    
    return {"error": f"No form found for intent: {intent_id}"}


@function_tool
def get_form_id_from_intent(intent_id: str) -> Dict[str, Any]:
    return get_form_id_from_intent_logic(intent_id)


def load_form_logic(form_id: str) -> Dict[str, Any]:
    """Load form definition from JSON config."""
    if form_id in _FORM_CACHE:
        return _FORM_CACHE[form_id]

    form_path = os.path.join("configs", "forms", f"{form_id}.json")
    if os.path.exists(form_path):
        with open(form_path, "r", encoding="utf-8") as f:
            _FORM_CACHE[form_id] = json.load(f)
            return _FORM_CACHE[form_id]
    return {"error": f"Form {form_id} not found"}


@function_tool
def load_form(form_id: str) -> Dict[str, Any]:
    return load_form_logic(form_id)


@function_tool
def get_next_form_field(form_id: str, session_id: str = "default", collected_fields_json: str = "{}") -> Dict[str, Any]:
    """Get the next field to ask for in form filling. One field at a time. collected_fields_json should be JSON string of collected fields."""
    form = load_form(form_id)
    if "error" in form:
        return form
    
    try:
        collected_fields = json.loads(collected_fields_json) if collected_fields_json else {}
    except Exception:
        collected_fields = {}
    
    # Check required fields first
    required_fields = form.get("required_fields", [])
    for field in sorted(required_fields, key=lambda x: x.get("order", 999)):
        field_id = field["field_id"]
        if field_id not in collected_fields:
            return {
                "field_id": field_id,
                "field_name": field["field_name"],
                "field_type": field["field_type"],
                "prompt": field["prompt"],
                "validation": field.get("validation"),
                "is_required": True,
                "next_field": True
            }
    
    # Then optional fields
    optional_fields = form.get("optional_fields", [])
    for field in sorted(optional_fields, key=lambda x: x.get("order", 999)):
        field_id = field["field_id"]
        if field_id not in collected_fields:
            return {
                "field_id": field_id,
                "field_name": field["field_name"],
                "field_type": field["field_type"],
                "prompt": field["prompt"],
                "validation": field.get("validation"),
                "is_required": False,
                "next_field": True
            }
    
    return {
        "next_field": False,
        "all_fields_collected": True,
        "message": "All fields collected"
    }


def validate_form_field_logic(field_id: str, field_value: str, validation_type: str) -> Dict[str, Any]:
    """Validate a single form field based on validation type."""
    if validation_type == "not_empty":
        if not field_value or not field_value.strip():
            return {"valid": False, "error": f"{field_id} cannot be empty"}
        return {"valid": True}
    
    elif validation_type == "email_format":
        import re
        email_pattern = r'^[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}$'
        if re.match(email_pattern, field_value):
            return {"valid": True}
        return {"valid": False, "error": "Invalid email format"}
    
    elif validation_type == "uk_phone":
        import re
        cleaned = field_value.replace(" ", "").replace("-", "")
        if re.match(r"^\+?[0-9().]{7,20}$", cleaned):
            return {"valid": True, "parsed_value": field_value.strip()}
        return {"valid": False, "error": "Please enter a valid phone number"}
    
    elif validation_type == "date_format":
        result = normalize_human_date_logic(field_value)
        if result.get("valid"):
            return {
                "valid": True,
                "parsed_date": result.get("full_date"),
                "parsed_value": result.get("full_date"),
                "iso_date": result.get("iso_date"),
            }
        return {"valid": False, "error": result.get("error", "Please enter the date in your own words")}
    
    elif validation_type == "time_format":
        result = normalize_human_time_logic(field_value)
        if result.get("valid"):
            return {
                "valid": True,
                "parsed_time": result.get("time_24h"),
                "parsed_value": result.get("display_time") or result.get("time_24h"),
            }
        return {"valid": False, "error": result.get("error", "Please enter the time in your own words")}
    
    return {"valid": True}  # Default: valid if no validation specified


@function_tool
def validate_form_field(field_id: str, field_value: str, validation_type: str) -> Dict[str, Any]:
    return validate_form_field_logic(field_id, field_value, validation_type)


@function_tool
def validate_complete_form(form_id: str, collected_fields_json: str) -> Dict[str, Any]:
    """Validate complete form: check missing fields and validate formats. collected_fields_json should be JSON string of collected fields."""
    form = load_form(form_id)
    if "error" in form:
        return form
    
    try:
        collected_fields = json.loads(collected_fields_json) if collected_fields_json else {}
    except Exception:
        return {"valid": False, "error": "Invalid JSON in collected_fields_json"}
    
    missing_fields = []
    validation_errors = []
    
    # Check required fields
    required_fields = form.get("required_fields", [])
    for field in required_fields:
        field_id = field["field_id"]
        if field_id not in collected_fields or not collected_fields[field_id]:
            missing_fields.append(field["field_name"])
        else:
            # Validate format
            validation_type = field.get("validation")
            if validation_type:
                validation_result = validate_form_field(field_id, collected_fields[field_id], validation_type)
                if not validation_result.get("valid"):
                    validation_errors.append(f"{field['field_name']}: {validation_result.get('error', 'Invalid')}")
    
    # Validate optional fields if provided
    optional_fields = form.get("optional_fields", [])
    for field in optional_fields:
        field_id = field["field_id"]
        if field_id in collected_fields and collected_fields[field_id]:
            validation_type = field.get("validation")
            if validation_type:
                validation_result = validate_form_field(field_id, collected_fields[field_id], validation_type)
                if not validation_result.get("valid"):
                    validation_errors.append(f"{field['field_name']}: {validation_result.get('error', 'Invalid')}")
    
    if missing_fields:
        return {
            "valid": False,
            "missing_fields": missing_fields,
            "message": f"Missing required fields: {', '.join(missing_fields)}"
        }
    
    if validation_errors:
        return {
            "valid": False,
            "validation_errors": validation_errors,
            "message": f"Validation errors: {', '.join(validation_errors)}"
        }
    
    return {
        "valid": True,
        "message": "Form is valid and complete",
        "form_data": collected_fields
    }


@function_tool
def generate_form_confirmation(form_id: str, collected_fields_json: str) -> str:
    """Generate confirmation message from form template. collected_fields_json should be JSON string of collected fields."""
    form = load_form(form_id)
    if "error" in form:
        return f"Error: {form['error']}"
    
    try:
        collected_fields = json.loads(collected_fields_json) if collected_fields_json else {}
    except Exception:
        return "Error: Invalid JSON in collected_fields_json"
    
    template = form.get("confirmation_template", "Please confirm:\n{fields}")
    
    # Simple template replacement
    confirmation = template
    for field_id, value in collected_fields.items():
        confirmation = confirmation.replace(f"{{{field_id}}}", str(value))
    
    return confirmation


@function_tool
def log_tool(message: str, agent: str):
    """Log events for agent trace."""
    log_event(agent, message)
    return True


# =================== BOOKING STATE MANAGEMENT ===================

def init_booking_logic(form_id: str) -> Dict[str, Any]:
    """Initialize a new booking session."""
    # Reset state for this session (simulated session scope)
    # In a real app, this would be keyed by session_id
    global _FORM_STATES
    session_id = "car_dealer_session_v2" # Aligned with Orchestrator session
    
    _FORM_STATES[session_id] = {
        "form_id": form_id,
        "collected": {},
        "current_field": None,
        "status": "in_progress", # in_progress, confirming, finalized
        "last_prompt": None
    }
    log_event("Booking", f"Initialized form: {form_id}")
    return {"success": True, "message": "Booking initialized"}


@function_tool
def init_booking(form_id: str) -> Dict[str, Any]:
    return init_booking_logic(form_id)


def get_booking_status_logic(session_id: str) -> str:
    """Get current booking status."""
    global _FORM_STATES
    state = _FORM_STATES.get(session_id)
    if not state:
        return "none"
    return state.get("status", "none")


@function_tool
def process_booking_step(user_input: str) -> Dict[str, Any]:
    return process_booking_step_logic(user_input)


def process_booking_step_logic(user_input: str) -> Dict[str, Any]:
    """
    Process one step of the booking flow.
    Handles validation, state update, and next question generation.
    Returns the exact message the agent should say.
    """
    global _FORM_STATES
    try:
        session_id = "car_dealer_session_v2"
        log_event("Booking", f"Processing step for session: {session_id}")
        state = _FORM_STATES.get(session_id)
        
        if not state:
            log_event("Booking", "State not found for session")
            return {"message": "Session expired. Please start over.", "status": "error"}
        
        form_id = state["form_id"]
        log_event("Booking", f"Loading form: {form_id}")
        form = load_form_logic(form_id)
        if "error" in form:
            log_event("Booking", f"Form load error: {form}")
            return {
                "message": (
                    f"I could not load the booking form '{form_id}'. "
                    "Please check that configs/forms contains the form definition."
                ),
                "status": "error",
            }
    
        collected = state["collected"]
        current_field_id = state["current_field"]
        status = state["status"]
        
        # 1. Handle Confirmation Phase
        if status == "confirming":
            if any(word in user_input.lower() for word in ["yes", "confirm", "correct", "ok", "sure"]):
                state["status"] = "finalized"
                
                # Generate final JSON output
                summary = "Great! Here's your booking:\n"
                for field_id, value in collected.items():
                    summary += f"- {field_id}: {value}\n"
                
                summary += "\n```json\n" + json.dumps(collected, indent=2) + "\n```\n"
                summary += "\nA team member will contact you shortly to confirm."
                
                return {
                    "message": summary,
                    "status": "finalized",
                    "data": collected
                }
            elif any(word in user_input.lower() for word in ["no", "wrong", "change"]):
                # Reset to in_progress to correct fields
                # For simplicity, we just ask what to change, but here we'll just restart the loop
                # A better implementation would parse "change date"
                state["status"] = "in_progress"
                return {
                    "message": "What would you like to change? (e.g., 'change date' or just provide the correct value)",
                    "status": "in_progress",
                    "data": collected
                }
            else:
                 return {
                    "message": "Please confirm if these details are correct (Yes/No).",
                    "status": "confirming",
                    "data": collected
                }
    
        # 2. Handle Input for Current Field
        if current_field_id:
            # Validate input
            # Find field definition
            field_def = next((f for f in form.get("required_fields", []) + form.get("optional_fields", []) if f["field_id"] == current_field_id), None)
            
            if field_def:
                validation_type = field_def.get("validation")
                # Special handling for "change X" commands could go here
                
                valid_res = validate_form_field_logic(current_field_id, user_input, validation_type)
                
                if valid_res.get("valid"):
                    # Save value
                    # Use normalized values when available (e.g. dates/times).
                    value_to_save = valid_res.get("parsed_value") or valid_res.get("parsed_date") or user_input
                    collected[current_field_id] = value_to_save
                    state["collected"] = collected
                    state["current_field"] = None # Reset so we pick next field
                else:
                    # Invalid
                    error_msg = valid_res.get("error", "Invalid input")
                    return {
                        "message": f"{error_msg}. {field_def['prompt']}",
                        "status": "in_progress",
                        "data": collected
                    }
    
        # 3. Pick Next Field
        # Check required fields
        for field in sorted(form.get("required_fields", []), key=lambda x: x.get("order", 999)):
            if field["field_id"] not in collected:
                state["current_field"] = field["field_id"]
                result = {
                    "message": field["prompt"],
                    "status": "in_progress",
                    "data": collected
                }
                log_event("Booking", f"Returning: {result}")
                return result
                
        # Check optional fields (simplified: just ask all for now, or skip if not configured)
        # In this simplified version, we assume optional fields are skipped unless we have logic to ask "do you want to add X?"
        # For now, let's just finish if required are done.
        
        # 4. All Done -> Confirmation
        state["status"] = "confirming"
        
        # Generate confirmation summary
        summary = "Please confirm these details are correct:\n"
        for field in form.get("required_fields", []):
            val = collected.get(field["field_id"], "N/A")
            summary += f"- {field['field_name']}: {val}\n"
            
        result = {
            "message": summary,
            "status": "confirming",
            "data": collected
        }
        log_event("Booking", f"Returning: {result}")
        return result
    except Exception as e:
        import traceback
        error_msg = f"Error in process_booking_step: {str(e)}\n{traceback.format_exc()}"
        log_event("Booking", error_msg)
        return {"message": "Technical error processing booking.", "status": "error"}


def run_handoff_logic(user_input: str, session_id: str = "car_dealer_session_v2") -> str:
    """Deterministic handoff flow without an extra LLM round-trip."""
    state = _HANDOFF_STATE.get(session_id)

    if not state:
        availability = check_agent_availability_logic()
        _HANDOFF_STATE[session_id] = {"awaiting_details": True}
        log_event("Handoff", f"Availability checked: {availability['available']}")
        if availability["available"]:
            return "An agent is available. Please provide your name and phone number so we can connect you."
        return "Our agents are currently offline. Please leave your name and phone number and we'll contact you."

    entities = extract_entities_logic(user_input).get("entities", {})
    phone = entities.get("phone")
    email = entities.get("email")
    name = user_input.strip()
    if phone:
        name = name.replace(phone, "").strip(" ,-")
    if email:
        name = name.replace(email, "").strip(" ,-")

    if phone or (name and len(name.split()) >= 1 and len(name) > 2):
        capture_lead_logic(
            customer_name=name or "Customer",
            email=email,
            phone=phone,
            interest_type="human_handoff",
        )
        _HANDOFF_STATE.pop(session_id, None)
        return "Thank you. Your details have been passed to our team."

    return "Please provide your name and phone number."


def run_booking_logic(user_input: str) -> str:
    """Run booking form step directly and return the user-facing message."""
    result = process_booking_step_logic(user_input)
    return result.get("message", "Sorry, I couldn't process that booking step.")


# =================== AGENT FACTORY (6 Core Agents) ===================

def build_agents(orgName: str = "Demo Motors", tz: str = "Europe/London", language: str = "en", dealer_name: str = None):
    dealerInfo = "Sales, Service, Parts, Finance"
    brands = "BMW, Audi, Mercedes, Toyota"
    locations = "London HQ; Service West; Parts North"
    alias = "AI Assistant"
    lang_name = languages.get(language, "English")

    # Get current date/time for Temporal Agent context
    now_iso = datetime.utcnow().isoformat() + "Z"
    try:
        import pytz
        now_dt = datetime.now(pytz.timezone(tz))
        now_formatted = now_dt.strftime("%A, %d %B %Y %H:%M")
        day_of_week = now_dt.strftime("%A")
        date_only = now_dt.strftime("%d %B %Y")
    except Exception:
        now_dt = datetime.utcnow()
        now_formatted = now_dt.strftime("%A, %d %B %Y %H:%M (UTC)")
        day_of_week = now_dt.strftime("%A")
        date_only = now_dt.strftime("%d %B %Y")

    if dealer_name:
        greeting_name = dealer_name
    elif orgName and orgName.strip() and orgName.lower().strip() != "demo motors":
        greeting_name = orgName.strip()
    else:
        greeting_name = "our dealership"

    # 1. Intent Agent (Lightweight Classification)
    intent_agent = Agent(
        name="Intent Agent",
        instructions="""
You are an Intent Classification Agent.
Your ONLY job is to classify the user's intent into one of these categories:
- service_booking (book service, repair, MOT, maintenance)
- test_drive (test drive, try car)
- human_handoff (speak to human, agent, person)
- general_enquiry (anything else)

OUTPUT FORMAT:
You must return ONLY a JSON object. Do not add any text.
{
  "intent_id": "service_booking" | "test_drive" | "human_handoff" | "general_enquiry",
  "form_id": "service_booking" | "test_drive" | null
}
""",
        model_settings=ModelSettings(model=GEN_MODEL, temperature=0.0),
    )

    # 2. Knowledge Agent (Strict RAG)
    dealer_context = f" for {dealer_name}" if dealer_name else ""
    knowledge_agent = Agent(
        name="Knowledge Agent",
        instructions=f"""
You are a Knowledge Agent specialized in retrieving dealership-specific information{dealer_context}.
You can ONLY answer questions about the currently selected dealer: {dealer_name if dealer_name else "General KB"}.

MANDATORY PROCESS:
1. Call semantic_search(query) to search the KB.
2. If semantic_search returns error/blocked, use that exact message.
3. If results found, answer ONLY using the retrieved chunks.
4. If NO results found, respond EXACTLY: "This information is not available in the knowledge base."

STRICT RULES:
- NO hallucinations.
- NO guessing.
- NO answering about other dealers.
""",
        tools=[semantic_search, log_tool],
        model_settings=ModelSettings(model=GEN_MODEL, temperature=0.0, max_tokens=400),
    )

    # 3. Booking Agent (Merged Form + QA)
    booking_agent = Agent(
        name="Booking Agent",
        instructions=f"""
You are a Booking Agent.
Your job is to collect information from the user to complete a booking form.

MANDATORY PROCESS:
1. Call process_booking_step(user_message)
2. The tool will return a JSON object with:
   - "message": The response you MUST say to the user.
   - "status": "in_progress" or "complete".
   - "data": The collected data so far.

3. You MUST output the "message" from the tool EXACTLY.
   - Do NOT add your own greeting.
   - Do NOT add "Is there anything else?".
   - Just say what the tool tells you.

4. If status is "complete", the tool message will include the confirmation request.
   - Wait for user to say "yes" or "confirm".
   - Call process_booking_step again.
   - If confirmed, the tool will output the final JSON.

CRITICAL:
- Do NOT try to manage the form state yourself.
- ALWAYS delegate to process_booking_step.
- If the user asks a question unrelated to the form, say "I need to finish this booking first." and repeat the tool's question.

Address the user in {lang_name}.
""",
        tools=[process_booking_step, log_tool],
        model_settings=ModelSettings(model=GEN_MODEL, temperature=0.1),
    )

    # 5. Handoff Agent
    handoff_agent = Agent(
        name="Handoff Agent",
        instructions=f"""
You are a Handoff Agent.
Your job is to connect the user with a human agent.

MANDATORY PROCESS:
1. Call check_agent_availability().
2. If available:
   - Ask: "An agent is available. Please provide your name and phone number so we can connect you."
3. If NOT available:
   - Ask: "Our agents are currently offline. Please leave your name and phone number and we'll contact you."
4. Once user provides info:
   - Call capture_lead(name, phone).
   - Respond: "Thank you. Your details have been passed to our team."

STRICT RULES:
- Use the exact phrases above.
- Do not add extra chatter.
""",
        tools=[check_agent_availability, capture_lead, log_tool],
        model_settings=ModelSettings(model=GEN_MODEL, temperature=0.1),
    )

    return {
        "intent": intent_agent,
        "booking": booking_agent,
        "knowledge": knowledge_agent,
        "handoff": handoff_agent,
    }


# =================== OUTPUT VALIDATION GUARDRAIL ===================

def validate_output_for_dealer_isolation(response: str, dealer_name: str = None) -> str:
    """
    Output validation is intentionally disabled.

    Dealer isolation is enforced before retrieval by selecting the dealer-specific
    FAISS index. Post-generation blocking was too aggressive for dealer names like
    "Charles_Hurst" vs "Charles Hurst" and blocked valid KB answers.
    """
    return response


# =================== RUN HELPERS ===================

async def async_run_orchestrator(
    user_input: str,
    orgName: str = "Demo Motors",
    tz: str = "Europe/London",
    language: str = "en",
    dealer_name: str = None,
    verbose: bool = True,
):
    AGENT_LOG.clear()
    session_id = "car_dealer_session_v2"
    
    # Layer 1: Set dealer context and pre-load KB
    if dealer_name:
        set_dealer_name(dealer_name)
        _load_kb_to_cache(dealer_name)
    else:
        set_dealer_name(None)
        _load_kb_to_cache(None)
    
    # 1. Dealer Mismatch Check (Rule-based Layer 1)
    mismatch = check_dealer_mismatch_logic(user_input)
    if mismatch["is_mismatch"]:
        log_event("Orchestrator", f"Blocked mismatch: {mismatch['message']}")
        return mismatch["message"], list(AGENT_LOG)

    # 2. Greeting Detection (Rule-based)
    greeting_pattern = r"\b(hi|hello|hey|good morning|good afternoon|good evening)\b"
    if re.search(greeting_pattern, user_input.lower()):
        dealer_display = dealer_name if dealer_name else "our dealership"
        return f"Hi there! I'm your {dealer_display} assistant. How can I help you today?", list(AGENT_LOG)

    log_event("System", f"Starting orchestrator run (dealer: {dealer_name})")

    current_status = get_booking_status_logic(session_id)
    is_booking_active = current_status in ["in_progress", "confirming"]
    intent_data = classify_intent_logic(user_input)
    intent_id = intent_data.get("intent_id", "general_enquiry")
    form_id = intent_data.get("form")
    log_event("Orchestrator", f"Classified intent: {intent_id}, form_id: {form_id}")

    if is_booking_active:
        if intent_id == "human_handoff":
            final = run_handoff_logic(user_input, session_id=session_id)
        else:
            final = run_booking_logic(user_input)
    elif intent_id in ["service_booking", "test_drive"]:
        if not form_id:
            form_res = get_form_id_from_intent_logic(intent_id)
            form_id = form_res.get("form_id", "service_booking")
        _HANDOFF_STATE.pop(session_id, None)
        init_booking_logic(form_id)
        final = run_booking_logic(user_input)
    elif intent_id == "human_handoff":
        final = run_handoff_logic(user_input, session_id=session_id)
    else:
        final = answer_knowledge_query_logic(user_input, dealer_name=dealer_name)

    final_validated = validate_output_for_dealer_isolation(final, dealer_name)
    log_event("System", "Orchestration complete")
    return final_validated, list(AGENT_LOG)


def run_orchestrator(user_input: str, **kwargs):
    return asyncio.run(async_run_orchestrator(user_input, **kwargs))


# =================== CLI MAIN ===================

async def main_async(args):
    if args.rebuild:
        if args.dealer:
            build_kb_from_data(dealer_name=args.dealer)
        else:
            build_kb_from_data()
        return

    if args.dealer:
        ensure_kb(dealer_name=args.dealer)
    else:
        ensure_kb()

    if args.query:
        ans, logs = await async_run_orchestrator(args.query, orgName=args.org, tz=args.tz, language=args.lang, dealer_name=args.dealer)
        print("\n=== FINAL ANSWER ===\n")
        print(ans)
        print("\n=== AGENT LOG ===")
        for row in logs:
            print(row)
        return

    print("\n🤖 Car Dealership Multi-agent (6 Agents)")
    print("Type 'exit' to quit.\n")
    while True:
        q = input("\n❓ User: ")
        if q.strip().lower() in ("exit", "quit"):
            break
        ans, logs = await async_run_orchestrator(q, orgName=args.org, tz=args.tz, language=args.lang, dealer_name=args.dealer)
        print("\n💬 Assistant:", ans)
        print("\n🔍 Agent Trace:")
        for row in logs:
            print(" ", row)


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--rebuild", action="store_true", help="Rebuild FAISS KB from ./Data")
    parser.add_argument("--query", type=str, default=None, help="Single query to run then exit")
    parser.add_argument("--org", type=str, default="Demo Motors")
    parser.add_argument("--tz", type=str, default="Europe/London")
    parser.add_argument("--lang", type=str, default="en")
    parser.add_argument("--dealer", type=str, default=None, help="Dealer name for dealer-specific RAG")
    args = parser.parse_args()
    asyncio.run(main_async(args))


if __name__ == "__main__":
    main()

